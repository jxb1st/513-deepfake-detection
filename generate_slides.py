"""Generate the 4-slide project proposal deck as a polished .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- design tokens ----------
PRIMARY = RGBColor(0x16, 0x36, 0x5C)   # deep navy
ACCENT  = RGBColor(0x2B, 0xB3, 0xA3)   # modern teal
TEXT    = RGBColor(0x2C, 0x2C, 0x2C)   # near-black
MUTED   = RGBColor(0x9A, 0xA5, 0xB1)   # cool gray
SOFT_BG = RGBColor(0xF3, 0xF6, 0xF8)   # very pale blue-gray
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
WATERMARK = RGBColor(0xEC, 0xF1, 0xF4) # extremely pale for big section numbers

HEAD_FONT = "Helvetica Neue"
BODY_FONT = "Helvetica Neue"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def textbox(slide, left, top, width, height, runs, *,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = list of (text, {size, bold, color, font, italic, space_after})."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    for i, item in enumerate(runs):
        text, style = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", align)
        if "space_after" in style:
            p.space_after = Pt(style["space_after"])
        if "space_before" in style:
            p.space_before = Pt(style["space_before"])
        run = p.add_run()
        run.text = text
        run.font.name = style.get("font", BODY_FONT)
        run.font.size = Pt(style.get("size", 14))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = style.get("color", TEXT)
    return box


def bullets(slide, left, top, width, height, items, *,
            size=13, color=TEXT, lead_color=ACCENT, gap=10, line_spacing=1.2):
    """items: list of strings (a sub-bullet starts with '> ')."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    for i, raw in enumerate(items):
        is_sub = raw.startswith("> ")
        text = raw[2:] if is_sub else raw
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap if not is_sub else gap // 2)
        p.line_spacing = line_spacing

        lead = p.add_run()
        if is_sub:
            lead.text = "      —  "
            lead.font.color.rgb = MUTED
        else:
            lead.text = "▍  "
            lead.font.color.rgb = lead_color
        lead.font.name = BODY_FONT
        lead.font.size = Pt(size)
        lead.font.bold = True

        body = p.add_run()
        body.text = text
        body.font.name = BODY_FONT
        body.font.size = Pt(size - (1 if is_sub else 0))
        body.font.color.rgb = MUTED if is_sub else color
    return box


def rect(slide, left, top, width, height, fill, *, line=None, line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def rounded(slide, left, top, width, height, fill, *, line=None, line_w=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def header(slide, section_num, eyebrow, title, *, watermark=True):
    """Standard slide header: big watermark number, eyebrow, title, accent rule."""
    if watermark:
        textbox(slide, Inches(10.4), Inches(0.05), Inches(2.9), Inches(2.2),
                [(section_num, {"size": 140, "bold": True, "color": WATERMARK,
                                "font": HEAD_FONT, "align": PP_ALIGN.RIGHT})])
    textbox(slide, Inches(0.7), Inches(0.55), Inches(8), Inches(0.35),
            [(eyebrow, {"size": 11, "bold": True, "color": ACCENT,
                        "font": HEAD_FONT})])
    textbox(slide, Inches(0.7), Inches(0.90), Inches(10), Inches(0.7),
            [(title, {"size": 30, "bold": True, "color": PRIMARY,
                      "font": HEAD_FONT})])
    rect(slide, Inches(0.7), Inches(1.55), Inches(0.6), Inches(0.045), ACCENT)


def footer(slide, page_num):
    textbox(slide, Inches(0.7), Inches(7.10), Inches(8), Inches(0.3),
            [("COURSE 513   •   PROJECT PROPOSAL   •   SPRING 2026",
              {"size": 9, "color": MUTED, "font": HEAD_FONT,
               "bold": True})])
    textbox(slide, Inches(11.8), Inches(7.10), Inches(1.0), Inches(0.3),
            [(f"{page_num}", {"size": 11, "bold": True, "color": ACCENT,
                              "font": HEAD_FONT, "align": PP_ALIGN.RIGHT}),
             (f"  /  4", {"size": 11, "color": MUTED,
                          "font": HEAD_FONT, "align": PP_ALIGN.RIGHT})])
    # Tiny corner accent
    rect(slide, Inches(0, ), Inches(7.42), Inches(13.333), Inches(0.08),
         SOFT_BG)


# ============================================================
# SLIDE 1 — Title
# ============================================================
s1 = prs.slides.add_slide(BLANK)

# Left navy column
rect(s1, 0, 0, Inches(5.0), Inches(7.5), PRIMARY)

# Small accent square on the left column
rect(s1, Inches(0.7), Inches(0.9), Inches(0.4), Inches(0.08), ACCENT)
textbox(s1, Inches(0.7), Inches(1.0), Inches(4), Inches(0.4),
        [("PROJECT PROPOSAL", {"size": 11, "bold": True,
                                "color": ACCENT, "font": HEAD_FONT})])
textbox(s1, Inches(0.7), Inches(1.4), Inches(4), Inches(0.35),
        [("Course 513  ·  Spring 2026",
          {"size": 11, "color": WHITE, "font": HEAD_FONT})])

# Author block (left bottom)
textbox(s1, Inches(0.7), Inches(6.0), Inches(4), Inches(0.3),
        [("PRESENTED BY", {"size": 10, "bold": True,
                            "color": ACCENT, "font": HEAD_FONT})])
textbox(s1, Inches(0.7), Inches(6.3), Inches(4), Inches(0.5),
        [("[ Your Name ]", {"size": 18, "bold": True,
                             "color": WHITE, "font": HEAD_FONT})])
textbox(s1, Inches(0.7), Inches(6.8), Inches(4), Inches(0.3),
        [("Introduction to Deep Learning",
          {"size": 10, "color": MUTED, "font": HEAD_FONT})])

# Right (white) area: big title
textbox(s1, Inches(5.5), Inches(1.6), Inches(7.5), Inches(0.4),
        [("A DEEP LEARNING PROPOSAL",
          {"size": 11, "bold": True, "color": ACCENT, "font": HEAD_FONT})])

textbox(s1, Inches(5.5), Inches(2.0), Inches(7.5), Inches(1.6),
        [("Real or Fake?", {"size": 64, "bold": True,
                              "color": PRIMARY, "font": HEAD_FONT})])

# Accent rule under title
rect(s1, Inches(5.5), Inches(3.55), Inches(1.0), Inches(0.06), ACCENT)

textbox(s1, Inches(5.5), Inches(3.8), Inches(7.5), Inches(1.6),
        [("A Deep Learning Approach to Detecting",
          {"size": 22, "color": TEXT, "font": HEAD_FONT,
           "space_after": 4}),
         ("AI-Generated Face Images",
          {"size": 22, "color": TEXT, "font": HEAD_FONT})])

# Decorative dots row (subtle, lower right)
for i, x in enumerate([5.5, 5.75, 6.0, 6.25, 6.5]):
    d = s1.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(x), Inches(6.45), Inches(0.10), Inches(0.10))
    d.fill.solid()
    d.fill.fore_color.rgb = ACCENT if i % 2 == 0 else MUTED
    d.line.fill.background()

textbox(s1, Inches(5.5), Inches(6.8), Inches(7.5), Inches(0.3),
        [("Computer Vision  ·  Image Classification  ·  Transfer Learning",
          {"size": 10, "bold": True, "color": MUTED, "font": HEAD_FONT})])


# ============================================================
# SLIDE 2 — Problem & Motivation
# ============================================================
s2 = prs.slides.add_slide(BLANK)
header(s2, "02", "PROBLEM  ·  MOTIVATION",
       "Why Detecting Deepfakes Matters in 2026")

bullets(s2, Inches(0.7), Inches(1.95), Inches(12.0), Inches(3.6), [
    "The problem  —  Generative models (StyleGAN3, Stable Diffusion, "
    "Midjourney v7) now produce photorealistic faces that fool human "
    "viewers more than 50% of the time.",
    "Real-world harm  —  identity fraud, fabricated evidence, election "
    "misinformation, non-consensual imagery, and fake social profiles at scale.",
    "The need  —  platforms, journalists, and courts require automated, "
    "scalable authenticity checks. Manual review cannot keep up.",
    "The opportunity  —  CNNs and Vision Transformers can detect sub-pixel "
    "artifacts (GAN fingerprints, frequency inconsistencies) invisible to people.",
], size=14, gap=12, line_spacing=1.25)

# Goal callout — soft background card, no border
rounded(s2, Inches(0.7), Inches(5.65), Inches(12.0), Inches(1.25), SOFT_BG)
# Accent left bar inside the card
rect(s2, Inches(0.7), Inches(5.65), Inches(0.10), Inches(1.25), ACCENT)

textbox(s2, Inches(1.0), Inches(5.78), Inches(11.5), Inches(0.3),
        [("PROJECT GOAL", {"size": 10, "bold": True,
                             "color": ACCENT, "font": HEAD_FONT})])
textbox(s2, Inches(1.0), Inches(6.10), Inches(11.5), Inches(0.75),
        [("Build a binary classifier that flags whether a face image is real "
          "or AI-generated — and honestly evaluate how well it generalizes "
          "beyond its training distribution.",
          {"size": 14, "color": PRIMARY, "font": HEAD_FONT})])

footer(s2, 2)


# ============================================================
# SLIDE 3 — Data & Method
# ============================================================
s3 = prs.slides.add_slide(BLANK)
header(s3, "03", "DATA  ·  APPROACH", "Data & Method")

# ---- Left: DATA ----
textbox(s3, Inches(0.7), Inches(1.95), Inches(6.0), Inches(0.35),
        [("DATA", {"size": 11, "bold": True,
                    "color": ACCENT, "font": HEAD_FONT})])
rect(s3, Inches(0.7), Inches(2.30), Inches(0.30), Inches(0.035), ACCENT)

bullets(s3, Inches(0.7), Inches(2.45), Inches(6.0), Inches(4.5), [
    "Primary dataset  —  \"140K Real and Fake Faces\" (Kaggle, ~3 GB)",
    "> 70,000 real faces from FFHQ (Flickr-Faces-HQ)",
    "> 70,000 fake faces generated by StyleGAN",
    "> 256×256 RGB, balanced binary labels",
    "> Pre-split into train / val / test",
    "Cross-generator test set  —  held-out batch of Stable-Diffusion and "
    "Midjourney faces, for generalization analysis.",
    "Preprocessing  —  face-crop check, resize to 224×224, ImageNet "
    "normalization, augmentation (flip, color jitter, JPEG simulation).",
], size=12, gap=8, line_spacing=1.22)

# Subtle vertical divider
rect(s3, Inches(6.9), Inches(2.0), Inches(0.015), Inches(4.9), MUTED)

# ---- Right: METHOD ----
textbox(s3, Inches(7.2), Inches(1.95), Inches(5.8), Inches(0.35),
        [("METHOD", {"size": 11, "bold": True,
                      "color": ACCENT, "font": HEAD_FONT})])
rect(s3, Inches(7.2), Inches(2.30), Inches(0.30), Inches(0.035), ACCENT)

bullets(s3, Inches(7.2), Inches(2.45), Inches(5.8), Inches(4.5), [
    "Model  —  transfer learning: fine-tune EfficientNet-B0 (primary); "
    "compare against ResNet-50 and ViT-B/16 (ImageNet-pretrained, timm).",
    "I/O  —  image tensor (3, 224, 224) → scalar probability (real vs fake).",
    "Loss  —  Binary Cross-Entropy with logits.",
    "Baselines  —  (1) shallow CNN from scratch; (2) FFT-magnitude features "
    "+ logistic regression (tests the frequency-artifact hypothesis).",
    "Metrics  —  Accuracy, AUC-ROC, F1, confusion matrix.",
    "Target  —  > 95% in-distribution accuracy; cross-generator AUC "
    "reported honestly.",
    "Feasibility  —  fits Colab T4. EfficientNet-B0 ≈ 1.5 hr/epoch; "
    "5 epochs is plenty.",
], size=12, gap=8, line_spacing=1.22)

footer(s3, 3)


# ============================================================
# SLIDE 4 — Limitations & Timeline
# ============================================================
s4 = prs.slides.add_slide(BLANK)
header(s4, "04", "RISKS  ·  PLAN", "What Could Go Wrong & Project Plan")

# ---- Left: LIMITATIONS ----
textbox(s4, Inches(0.7), Inches(1.95), Inches(6.0), Inches(0.35),
        [("LIMITATIONS, BIASES & CONCERNS",
          {"size": 11, "bold": True, "color": ACCENT, "font": HEAD_FONT})])
rect(s4, Inches(0.7), Inches(2.30), Inches(0.30), Inches(0.035), ACCENT)

bullets(s4, Inches(0.7), Inches(2.45), Inches(6.0), Inches(4.5), [
    "Generalization gap  —  trained on StyleGAN; may fail on Diffusion "
    "images. The central honesty question of this project.",
    "Demographic bias  —  FFHQ skews younger, lighter-skinned, Western. "
    "Will measure performance per demographic slice where possible.",
    "Adversarial fragility  —  tiny perturbations or JPEG re-compression "
    "can flip predictions (well documented in the literature).",
    "Ethical risk  —  false positives can wrongly accuse genuine photos. "
    "Will report calibration, not just accuracy.",
    "Dataset staleness  —  StyleGAN is 2019; 2026 threat is mostly "
    "Diffusion. Mitigated via cross-generator evaluation.",
], size=12, gap=9, line_spacing=1.22)

# Divider
rect(s4, Inches(6.9), Inches(2.0), Inches(0.015), Inches(4.9), MUTED)

# ---- Right: TIMELINE ----
textbox(s4, Inches(7.2), Inches(1.95), Inches(5.8), Inches(0.35),
        [("8-WEEK TIMELINE", {"size": 11, "bold": True,
                                "color": ACCENT, "font": HEAD_FONT})])
rect(s4, Inches(7.2), Inches(2.30), Inches(0.30), Inches(0.035), ACCENT)

# Visual timeline rows
timeline = [
    ("Weeks 1–2", "Data download, EDA, preprocessing pipeline"),
    ("Weeks 3–4", "Baseline (FFT + LR) + first fine-tuned CNN"),
    ("Weeks 5–6", "Model comparison (ResNet / EfficientNet / ViT) + tuning"),
    ("Week 7",    "Cross-generator generalization + bias analysis"),
    ("Week 8",    "Final report, slides, code cleanup"),
]
row_top = 2.55
for idx, (label, desc) in enumerate(timeline):
    y = row_top + idx * 0.55
    # accent dot
    dot = s4.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(7.2), Inches(y + 0.07),
                               Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    # label
    textbox(s4, Inches(7.45), Inches(y), Inches(1.4), Inches(0.4),
            [(label, {"size": 11, "bold": True,
                      "color": PRIMARY, "font": HEAD_FONT})])
    # description
    textbox(s4, Inches(8.85), Inches(y), Inches(4.2), Inches(0.4),
            [(desc, {"size": 11, "color": TEXT, "font": HEAD_FONT})])

# Stretch goal card
rounded(s4, Inches(7.2), Inches(5.55), Inches(5.8), Inches(1.35), SOFT_BG)
rect(s4, Inches(7.2), Inches(5.55), Inches(0.10), Inches(1.35), ACCENT)
textbox(s4, Inches(7.45), Inches(5.68), Inches(5.4), Inches(0.3),
        [("STRETCH GOAL", {"size": 10, "bold": True,
                             "color": ACCENT, "font": HEAD_FONT})])
textbox(s4, Inches(7.45), Inches(5.98), Inches(5.4), Inches(0.85),
        [("Use Grad-CAM to visualize what the model attends to — eyes, "
          "hairline, or background? Strong narrative payoff for the final talk.",
          {"size": 11, "color": PRIMARY, "font": HEAD_FONT})])

footer(s4, 4)


out_path = "/Users/jxb1st/Desktop/513project/Project_Proposal.pptx"
prs.save(out_path)
print(f"Saved {out_path}")
